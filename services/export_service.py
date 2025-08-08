import os
import json
import markdown
from datetime import datetime
from typing import Dict, Any, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from config import Config
import uuid

class ExportService:
    def __init__(self):
        self.export_folder = Config.EXPORT_FOLDER
        os.makedirs(self.export_folder, exist_ok=True)
    
    def to_markdown(self, project, options: Dict[str, Any] = None) -> Tuple[str, str]:
        """
        Generate Markdown export from project data
        """
        if not options:
            options = {}
        
        filename = f"{project.title or 'whiteboard'}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
        file_path = os.path.join(self.export_folder, filename)
        
        try:
            # Collect all structured content from whiteboards
            all_content = []
            for whiteboard in project.whiteboards:
                if whiteboard.processing_status == 'completed' and whiteboard.structured_content:
                    content = whiteboard.get_structured_content()
                    if content:
                        all_content.append(content)
            
            if not all_content:
                raise ValueError("No processed content available for export")
            
            # Generate markdown content
            md_content = self._generate_markdown_content(all_content, project, options)
            
            # Write to file
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(md_content)
            
            return file_path, filename
        
        except Exception as e:
            raise Exception(f"Markdown export failed: {str(e)}")
    
    def to_pptx(self, project, options: Dict[str, Any] = None) -> Tuple[str, str]:
        """
        Generate PowerPoint presentation from project data
        """
        if not options:
            options = {}
        
        filename = f"{project.title or 'whiteboard'}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path = os.path.join(self.export_folder, filename)
        
        try:
            # Collect all structured content
            all_content = []
            for whiteboard in project.whiteboards:
                if whiteboard.processing_status == 'completed' and whiteboard.structured_content:
                    content = whiteboard.get_structured_content()
                    if content:
                        all_content.append(content)
            
            if not all_content:
                raise ValueError("No processed content available for export")
            
            # Create presentation
            prs = Presentation()
            self._create_powerpoint_slides(prs, all_content, project, options)
            
            # Save presentation
            prs.save(file_path)
            
            return file_path, filename
        
        except Exception as e:
            raise Exception(f"PowerPoint export failed: {str(e)}")
    
    def to_mindmap(self, project, options: Dict[str, Any] = None) -> Tuple[str, str]:
        """
        Generate mind map data from project content
        """
        if not options:
            options = {}
        
        format_type = options.get('format', 'json')
        filename = f"{project.title or 'whiteboard'}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{format_type}"
        file_path = os.path.join(self.export_folder, filename)
        
        try:
            # Collect all structured content
            all_content = []
            for whiteboard in project.whiteboards:
                if whiteboard.processing_status == 'completed' and whiteboard.structured_content:
                    content = whiteboard.get_structured_content()
                    if content:
                        all_content.append(content)
            
            if not all_content:
                raise ValueError("No processed content available for export")
            
            # Generate mind map structure
            mindmap_data = self._generate_mindmap_structure(all_content, project, options)
            
            # Export in requested format
            if format_type == 'json':
                with open(file_path, 'w', encoding='utf-8') as f:
                    json.dump(mindmap_data, f, indent=2, ensure_ascii=False)
            elif format_type == 'xmind':
                # For demo, we'll save as JSON with xmind structure
                xmind_data = self._convert_to_xmind_format(mindmap_data)
                with open(file_path, 'w', encoding='utf-8') as f:
                    json.dump(xmind_data, f, indent=2, ensure_ascii=False)
            elif format_type == 'freemind':
                # For demo, we'll save as XML structure
                freemind_xml = self._convert_to_freemind_format(mindmap_data)
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write(freemind_xml)
            
            return file_path, filename
        
        except Exception as e:
            raise Exception(f"Mind map export failed: {str(e)}")
    
    def to_notion(self, project, options: Dict[str, Any] = None) -> Tuple[str, str]:
        """
        Generate Notion-compatible export
        """
        if not options:
            options = {}
        
        filename = f"{project.title or 'whiteboard'}_notion_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"
        file_path = os.path.join(self.export_folder, filename)
        
        try:
            # Collect all structured content
            all_content = []
            for whiteboard in project.whiteboards:
                if whiteboard.processing_status == 'completed' and whiteboard.structured_content:
                    content = whiteboard.get_structured_content()
                    if content:
                        all_content.append(content)
            
            if not all_content:
                raise ValueError("No processed content available for export")
            
            # Generate Notion-specific markdown
            notion_content = self._generate_notion_content(all_content, project, options)
            
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(notion_content)
            
            return file_path, filename
        
        except Exception as e:
            raise Exception(f"Notion export failed: {str(e)}")
    
    def to_confluence(self, project, options: Dict[str, Any] = None) -> Tuple[str, str]:
        """
        Generate Confluence wiki format export
        """
        if not options:
            options = {}
        
        filename = f"{project.title or 'whiteboard'}_confluence_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt"
        file_path = os.path.join(self.export_folder, filename)
        
        try:
            # Collect all structured content
            all_content = []
            for whiteboard in project.whiteboards:
                if whiteboard.processing_status == 'completed' and whiteboard.structured_content:
                    content = whiteboard.get_structured_content()
                    if content:
                        all_content.append(content)
            
            if not all_content:
                raise ValueError("No processed content available for export")
            
            # Generate Confluence markup
            confluence_content = self._generate_confluence_content(all_content, project, options)
            
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(confluence_content)
            
            return file_path, filename
        
        except Exception as e:
            raise Exception(f"Confluence export failed: {str(e)}")
    
    def _generate_markdown_content(self, all_content: list, project, options: Dict) -> str:
        """
        Generate markdown content from structured data
        """
        md_lines = []
        
        # Title
        title = project.title or "Meeting Whiteboard Notes"
        md_lines.append(f"# {title}")
        md_lines.append("")
        
        # Metadata
        md_lines.append("## Meeting Information")
        md_lines.append(f"- **Date**: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        md_lines.append(f"- **Project**: {project.title or 'Untitled'}")
        if project.description:
            md_lines.append(f"- **Description**: {project.description}")
        md_lines.append("")
        
        # Process each whiteboard's content
        for i, content in enumerate(all_content):
            if len(all_content) > 1:
                md_lines.append(f"## Whiteboard {i + 1}")
                md_lines.append("")
            
            # Sections
            if content.get('sections'):
                for section in content['sections']:
                    md_lines.append(f"### {section.get('heading', 'Section')}")
                    md_lines.append("")
                    if section.get('content'):
                        md_lines.append(section['content'])
                    md_lines.append("")
                    
                    # Subsections
                    if section.get('subsections'):
                        for subsection in section['subsections']:
                            md_lines.append(f"#### {subsection.get('heading', 'Subsection')}")
                            md_lines.append("")
                            if subsection.get('content'):
                                md_lines.append(subsection['content'])
                            md_lines.append("")
            
            # Tables
            if content.get('tables'):
                md_lines.append("### Tables")
                md_lines.append("")
                for table in content['tables']:
                    if table.get('title'):
                        md_lines.append(f"#### {table['title']}")
                        md_lines.append("")
                    
                    headers = table.get('headers', [])
                    rows = table.get('rows', [])
                    
                    if headers:
                        md_lines.append("| " + " | ".join(headers) + " |")
                        md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                        
                        for row in rows:
                            md_lines.append("| " + " | ".join(str(cell) for cell in row) + " |")
                        md_lines.append("")
            
            # Action Items
            if content.get('action_items'):
                md_lines.append("### Action Items")
                md_lines.append("")
                for item in content['action_items']:
                    priority = item.get('priority', 'medium')
                    task = item.get('task', '')
                    assignee = item.get('assignee', '')
                    
                    priority_emoji = {"high": "🔴", "medium": "🟡", "low": "🟢"}.get(priority, "⚪")
                    assignee_text = f" (@{assignee})" if assignee else ""
                    
                    md_lines.append(f"- [ ] {priority_emoji} {task}{assignee_text}")
                md_lines.append("")
            
            # Key Points
            if content.get('key_points'):
                md_lines.append("### Key Points")
                md_lines.append("")
                for point in content['key_points']:
                    md_lines.append(f"- {point}")
                md_lines.append("")
            
            # Diagrams
            if content.get('diagrams'):
                md_lines.append("### Diagrams and Visuals")
                md_lines.append("")
                for diagram in content['diagrams']:
                    md_lines.append(f"#### {diagram.get('type', 'Diagram').title()}")
                    md_lines.append("")
                    md_lines.append(diagram.get('description', 'No description available'))
                    
                    if diagram.get('elements'):
                        md_lines.append("")
                        md_lines.append("**Elements:**")
                        for element in diagram['elements']:
                            md_lines.append(f"- {element}")
                    md_lines.append("")
        
        return "\n".join(md_lines)
    
    def _create_powerpoint_slides(self, prs: Presentation, all_content: list, project, options: Dict):
        """
        Create PowerPoint slides from content
        """
        # Title slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = project.title or "Meeting Whiteboard Notes"
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Process content slides
        for i, content in enumerate(all_content):
            # Overview slide for each whiteboard
            if len(all_content) > 1:
                slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = f"Whiteboard {i + 1}"
                
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.text = content.get('title', f'Whiteboard {i + 1} Content')
            
            # Sections slides
            if content.get('sections'):
                for section in content['sections']:
                    slide_layout = prs.slide_layouts[1]  # Title and content
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = section.get('heading', 'Section')
                    
                    content_shape = slide.placeholders[1]
                    tf = content_shape.text_frame
                    tf.text = section.get('content', '')[:500] + '...' if len(section.get('content', '')) > 500 else section.get('content', '')
            
            # Action items slide
            if content.get('action_items'):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = "Action Items"
                
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.clear()
                
                for item in content['action_items']:
                    p = tf.add_paragraph()
                    p.text = f"• {item.get('task', '')}"
                    if item.get('assignee'):
                        p.text += f" ({item['assignee']})"
                    
                    # Color code by priority
                    priority = item.get('priority', 'medium')
                    if priority == 'high':
                        p.font.color.rgb = RGBColor(220, 20, 60)  # Crimson
                    elif priority == 'medium':
                        p.font.color.rgb = RGBColor(255, 140, 0)   # Dark orange
                    else:
                        p.font.color.rgb = RGBColor(34, 139, 34)   # Forest green
            
            # Tables slide
            if content.get('tables'):
                for table_data in content['tables']:
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = table_data.get('title', 'Table')
                    
                    # Add table (simplified for demo)
                    headers = table_data.get('headers', [])
                    rows = table_data.get('rows', [])
                    
                    if headers and rows:
                        # Create a text representation
                        content_shape = slide.placeholders[1]
                        tf = content_shape.text_frame
                        tf.text = "Table Content:"
                        
                        for row in [headers] + rows:
                            p = tf.add_paragraph()
                            p.text = " | ".join(str(cell) for cell in row)
    
    def _generate_mindmap_structure(self, all_content: list, project, options: Dict) -> Dict:
        """
        Generate mind map structure
        """
        mindmap = {
            "meta": {
                "name": project.title or "Meeting Whiteboard",
                "created": datetime.now().isoformat(),
                "version": "1.0"
            },
            "root": {
                "topic": project.title or "Meeting Notes",
                "id": str(uuid.uuid4()),
                "children": []
            }
        }
        
        for i, content in enumerate(all_content):
            # Create main branch for each whiteboard
            whiteboard_branch = {
                "topic": f"Whiteboard {i + 1}" if len(all_content) > 1 else "Content",
                "id": str(uuid.uuid4()),
                "children": []
            }
            
            # Add sections as children
            if content.get('sections'):
                for section in content['sections']:
                    section_node = {
                        "topic": section.get('heading', 'Section'),
                        "id": str(uuid.uuid4()),
                        "children": []
                    }
                    
                    # Add subsections
                    if section.get('subsections'):
                        for subsection in section['subsections']:
                            subsection_node = {
                                "topic": subsection.get('heading', 'Subsection'),
                                "id": str(uuid.uuid4())
                            }
                            section_node['children'].append(subsection_node)
                    
                    whiteboard_branch['children'].append(section_node)
            
            # Add action items branch
            if content.get('action_items'):
                actions_branch = {
                    "topic": "Action Items",
                    "id": str(uuid.uuid4()),
                    "children": []
                }
                
                for item in content['action_items']:
                    action_node = {
                        "topic": item.get('task', 'Task')[:40] + '...' if len(item.get('task', '')) > 40 else item.get('task', ''),
                        "id": str(uuid.uuid4()),
                        "priority": item.get('priority', 'medium')
                    }
                    actions_branch['children'].append(action_node)
                
                whiteboard_branch['children'].append(actions_branch)
            
            # Add key points branch
            if content.get('key_points'):
                keypoints_branch = {
                    "topic": "Key Points",
                    "id": str(uuid.uuid4()),
                    "children": []
                }
                
                for point in content['key_points']:
                    point_node = {
                        "topic": point[:50] + '...' if len(point) > 50 else point,
                        "id": str(uuid.uuid4())
                    }
                    keypoints_branch['children'].append(point_node)
                
                whiteboard_branch['children'].append(keypoints_branch)
            
            mindmap['root']['children'].append(whiteboard_branch)
        
        return mindmap
    
    def _convert_to_xmind_format(self, mindmap_data: Dict) -> Dict:
        """
        Convert to XMind compatible format
        """
        return {
            "title": mindmap_data['meta']['name'],
            "structure": "org.xmind.ui.map.unbalanced",
            "theme": "xmind.ui.theme.default",
            "rootTopic": mindmap_data['root']
        }
    
    def _convert_to_freemind_format(self, mindmap_data: Dict) -> str:
        """
        Convert to FreeMind XML format
        """
        def node_to_xml(node, level=0):
            indent = "  " * level
            topic = node.get('topic', 'Node')
            xml = f'{indent}<node TEXT="{topic}">\n'
            
            for child in node.get('children', []):
                xml += node_to_xml(child, level + 1)
            
            xml += f'{indent}</node>\n'
            return xml
        
        xml_header = '<?xml version="1.0" encoding="UTF-8"?>\n<map version="1.0.1">\n'
        xml_footer = '</map>'
        
        return xml_header + node_to_xml(mindmap_data['root']) + xml_footer
    
    def _generate_notion_content(self, all_content: list, project, options: Dict) -> str:
        """
        Generate Notion-compatible markdown
        """
        lines = []
        
        # Page properties (Notion style)
        page_title = options.get('page_title', project.title or 'Meeting Notes')
        lines.append(f"# {page_title}")
        lines.append("")
        
        # Properties block
        if options.get('include_properties', True):
            lines.append("**Properties:**")
            lines.append(f"- Date: {datetime.now().strftime('%Y-%m-%d')}")
            lines.append(f"- Type: Meeting Notes")
            lines.append(f"- Status: Complete")
            lines.append("")
        
        # Content blocks
        for content in all_content:
            # Callout for key information
            if content.get('key_points'):
                lines.append("> 💡 **Key Takeaways**")
                for point in content['key_points'][:3]:  # Top 3 points
                    lines.append(f"> {point}")
                lines.append("")
            
            # Toggle blocks for sections
            if content.get('sections'):
                for section in content['sections']:
                    lines.append(f"## {section.get('heading', 'Section')}")
                    lines.append("")
                    if section.get('content'):
                        lines.append(section['content'])
                    lines.append("")
            
            # Database-style action items
            if content.get('action_items'):
                lines.append("## 📋 Action Items")
                lines.append("")
                lines.append("| Task | Priority | Assignee | Status |")
                lines.append("|------|----------|----------|--------|")
                
                for item in content['action_items']:
                    task = item.get('task', '')
                    priority = item.get('priority', 'medium')
                    assignee = item.get('assignee', 'Unassigned')
                    status = "To Do"
                    
                    lines.append(f"| {task} | {priority.title()} | {assignee} | {status} |")
                lines.append("")
        
        return "\n".join(lines)
    
    def _generate_confluence_content(self, all_content: list, project, options: Dict) -> str:
        """
        Generate Confluence wiki markup
        """
        lines = []
        
        # Page title
        lines.append(f"h1. {project.title or 'Meeting Whiteboard Notes'}")
        lines.append("")
        
        # Info macro
        lines.append("{info}")
        lines.append(f"*Generated:* {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        lines.append(f"*Project:* {project.title or 'Untitled'}")
        lines.append("{info}")
        lines.append("")
        
        for content in all_content:
            # Sections
            if content.get('sections'):
                for section in content['sections']:
                    lines.append(f"h2. {section.get('heading', 'Section')}")
                    lines.append("")
                    if section.get('content'):
                        lines.append(section['content'])
                    lines.append("")
            
            # Tables
            if content.get('tables'):
                for table in content['tables']:
                    if table.get('title'):
                        lines.append(f"h3. {table['title']}")
                        lines.append("")
                    
                    headers = table.get('headers', [])
                    rows = table.get('rows', [])
                    
                    if headers:
                        lines.append("|| " + " || ".join(headers) + " ||")
                        for row in rows:
                            lines.append("| " + " | ".join(str(cell) for cell in row) + " |")
                        lines.append("")
            
            # Action items with task list macro
            if content.get('action_items'):
                lines.append("h2. Action Items")
                lines.append("")
                lines.append("{task-list}")
                for item in content['action_items']:
                    assignee = item.get('assignee', '')
                    task = item.get('task', '')
                    if assignee:
                        lines.append(f"* {task} - [~{assignee}]")
                    else:
                        lines.append(f"* {task}")
                lines.append("{task-list}")
                lines.append("")
            
            # Key points in note macro
            if content.get('key_points'):
                lines.append("{note}")
                lines.append("h3. Key Points")
                for point in content['key_points']:
                    lines.append(f"* {point}")
                lines.append("{note}")
                lines.append("")
        
        return "\n".join(lines)